"""
IBM Z Newsletter – PPTX Builder (vollständig programmatisch)

Alle Folien werden von Grund auf gebaut.
Aus dem Template werden nur die IBM-Bilder (Banner, Logos) entnommen.

Struktur:
  Folie 1  : Cover  – linke Seitenleiste (THEMEN + EVENTS) + rechts Artikel 1
  Folie 2+ : Inhaltsfolien – 2 Artikel pro Folie
  Optional : Event-Folien – 3 Events pro Folie
  Letzte   : Abschlussfolie aus Template (unverändert)
"""

import copy, io, os, re, shutil
from typing import List, Optional

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

from config import OUTPUT_DIR, TEMPLATE_FILE, EVENTS_URL

# ── Farben ────────────────────────────────────────────────────────────────────
IBM_BLUE   = RGBColor(0x00, 0x43, 0xCE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x16, 0x16, 0x16)
MID_GRAY   = RGBColor(0x6F, 0x6F, 0x6F)
DIV_GRAY   = RGBColor(0xD8, 0xD8, 0xD8)
BAR_GRAY   = RGBColor(0xE7, 0xE6, 0xE6)

FONT = "IBM Plex Sans"

GERMAN_MONTHS = {
    1:"Januar", 2:"Februar", 3:"März",    4:"April",
    5:"Mai",    6:"Juni",    7:"Juli",     8:"August",
    9:"September", 10:"Oktober", 11:"November", 12:"Dezember",
}

# ── Geometrie (A4 Hochformat) ─────────────────────────────────────────────────
W, H          = Inches(8.27), Inches(11.69)
HEADER_H      = Inches(0.66)
FOOTER_Y      = Inches(11.10)
FOOTER_H      = Inches(0.65)

# Seitenleiste (Cover)
SIDE_W        = Inches(2.80)
SIDE_COL_L    = Inches(0.14)
SIDE_COL_W    = SIDE_W - SIDE_COL_L - Inches(0.10)

# Inhaltsbereich
CONTENT_TOP   = Inches(0.82)
CONTENT_BOT   = Inches(10.95)
MARGIN        = Inches(0.28)
NUM_W         = Inches(0.38)
TEXT_L_FULL   = MARGIN + NUM_W + Inches(0.06)
TEXT_W_FULL   = W - TEXT_L_FULL - Inches(0.22)

# Text-Spalte auf Cover (rechte Seite)
COVER_R_L     = SIDE_W + Inches(0.20)
COVER_R_W     = W - COVER_R_L - Inches(0.20)


# ═══════════════════════════════════════════════════════════════════════════════
# Hilfsfunktionen
# ═══════════════════════════════════════════════════════════════════════════════

def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def _tb(slide, l, t, w, h, text, size,
        bold=False, color=NEAR_BLACK, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = FONT
    return box


def _multi_para_tb(slide, l, t, w, h, text, size, color=NEAR_BLACK):
    """Textfeld mit Absätzen – getrennt durch Leerzeilen im Text."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    if not paragraphs:
        paragraphs = [text.strip()]
    for i, para_text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = Pt(5)
        r = p.add_run()
        r.text = para_text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = FONT
    return box


def _link_tb(slide, l, t, w, h, text, url):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9); r.font.color.rgb = IBM_BLUE
    r.font.underline = True; r.font.name = FONT
    try:
        rId = slide.part.relate_to(
            url,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True)
        rPr = r._r.get_or_add_rPr()
        hl = etree.SubElement(rPr, qn("a:hlinkClick"))
        hl.set(qn("r:id"), rId)
    except Exception:
        pass
    return box


def _divider(slide, y):
    _rect(slide, MARGIN, y, W - MARGIN * 2, Inches(0.015), DIV_GRAY)


def _picture(slide, img_bytes, l, t, w, h):
    if img_bytes:
        slide.shapes.add_picture(io.BytesIO(img_bytes), l, t, w, h)


def _get_image_bytes(article: dict):
    b = article.get("image_bytes")
    return b if b else None


# ── Bilder aus Template extrahieren ──────────────────────────────────────────

def _extract_images(prs):
    banner = logo_hd = None

    for shape in prs.slides[0].shapes:
        if shape.shape_type == 13 and shape.width > Inches(7.0):
            banner = shape.image.blob
            break

    for slide in list(prs.slides)[1:]:
        for shape in slide.shapes:
            if (shape.shape_type == 13
                    and shape.top < Inches(1.0)
                    and shape.left > Inches(6.0)):
                logo_hd = shape.image.blob
                break
        if logo_hd:
            break

    return banner, logo_hd


# ══════════════════════════════════════════════════════════════════════════════
# Header & Footer
# ══════════════════════════════════════════════════════════════════════════════

def _header(slide, logo_bytes):
    _rect(slide, Inches(-0.02), Inches(0), W + Inches(0.04), HEADER_H, BAR_GRAY)
    _picture(slide, logo_bytes, Inches(7.16), Inches(0.17), Inches(0.82), Inches(0.33))


def _footer(slide, logo_bytes, page_num, year):
    _rect(slide, Inches(-0.02), FOOTER_Y, W + Inches(0.04), FOOTER_H, BAR_GRAY)
    _tb(slide, Inches(0.10), Inches(11.44), Inches(2.20), Inches(0.22),
        f"© {year} IBM Corporation", 8, color=NEAR_BLACK)
    _tb(slide, Inches(3.90), Inches(11.36), Inches(0.50), Inches(0.28),
        str(page_num), 12, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    _picture(slide, logo_bytes, Inches(7.10), Inches(11.25), Inches(0.90), Inches(0.35))


# ══════════════════════════════════════════════════════════════════════════════
# Artikel-Block
# ══════════════════════════════════════════════════════════════════════════════

def _article_block(slide, article, num, top, block_h,
                   l_num, l_text, w_text, num_color=IBM_BLUE, text_color=NEAR_BLACK):
    pad = Inches(0.10)
    y   = top + pad

    _tb(slide, l_num, y, NUM_W, Inches(0.45),
        str(num), 16, bold=True, color=num_color, align=PP_ALIGN.CENTER, wrap=False)

    title_h = Inches(0.55)
    _tb(slide, l_text, y, w_text, title_h,
        article["title"], 11, bold=True, color=text_color)

    author_top = top + block_h - Inches(0.28)
    link_top   = author_top - Inches(0.30)

    pub = article.get("published")
    date_str = (f"{pub.day}. {GERMAN_MONTHS.get(pub.month,'')} {pub.year}"
                if pub else "")
    _link_tb(slide, l_text, link_top, w_text, Inches(0.26),
             "→ Mehr Informationen erhalten Sie hier",
             article.get("url", ""))
    _tb(slide, l_text, author_top, w_text, Inches(0.26),
        f"{article.get('author','')}  ·  {date_str}",
        8, color=MID_GRAY, italic=True)

    sum_top    = y + title_h + Inches(0.08)
    content_h  = link_top - sum_top - Inches(0.10)

    img_bytes = _get_image_bytes(article)
    MAX_IMG_H = min(Inches(1.60), content_h * 0.38)
    MAX_IMG_W = w_text

    if img_bytes:
        try:
            tmp = slide.shapes.add_picture(io.BytesIO(img_bytes), l_text, Inches(0))
            nat_w, nat_h = tmp.width, tmp.height
            tmp._element.getparent().remove(tmp._element)

            scale = min(1.0, MAX_IMG_W / nat_w, MAX_IMG_H / nat_h)
            img_w = int(nat_w * scale)
            img_h = int(nat_h * scale)

            sum_h   = max(content_h - img_h - Inches(0.12), Inches(0.80))
            img_top = sum_top + sum_h + Inches(0.10)
            img_left = l_text + Inches(0.05)
            slide.shapes.add_picture(io.BytesIO(img_bytes),
                                     img_left, img_top, img_w, img_h)
        except Exception:
            img_bytes = None

    if not img_bytes:
        sum_h = max(content_h, Inches(0.80))

    _multi_para_tb(slide, l_text, sum_top, w_text, sum_h,
                   article["summary"], 9, color=text_color)


# ══════════════════════════════════════════════════════════════════════════════
# Event-Block (für Event-Inhaltsfolien)
# ══════════════════════════════════════════════════════════════════════════════

def _event_block(slide, event, num, top, block_h,
                 l_num, l_text, w_text, num_color=IBM_BLUE, text_color=NEAR_BLACK):
    """Zeichnet einen Event-Block auf einer Inhaltsfolie."""
    pad = Inches(0.12)
    y   = top + pad

    # Datum statt Nummer (z.B. "15.\n04.")
    pub = event.get("event_date")
    if pub:
        date_label = f"{pub.day:02d}.\n{pub.month:02d}."
        _tb(slide, l_num, y, NUM_W, Inches(0.55),
            date_label, 8, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    else:
        _tb(slide, l_num, y, NUM_W, Inches(0.45),
            str(num), 16, bold=True, color=num_color, align=PP_ALIGN.CENTER)

    # Titel
    title_h = Inches(0.50)
    _tb(slide, l_text, y, w_text, title_h,
        event["title"], 11, bold=True, color=text_color)
    y += title_h + Inches(0.06)

    # Datum · Uhrzeit · Ort
    info_parts = []
    pub = event.get("event_date")
    if pub:
        info_parts.append(f"{pub.day}. {GERMAN_MONTHS.get(pub.month, '')} {pub.year}")
    if event.get("time_str"):
        info_parts.append(event["time_str"])
    if event.get("location"):
        info_parts.append(event["location"])

    if info_parts:
        _tb(slide, l_text, y, w_text, Inches(0.28),
            "  ·  ".join(info_parts), 9, color=MID_GRAY, italic=True)
        y += Inches(0.32)

    # Beschreibung
    desc = event.get("description", "")
    url  = event.get("url", "")
    link_y = top + block_h - Inches(0.28)

    if desc:
        desc_h = max(link_y - y - Inches(0.12), Inches(0.30))
        _multi_para_tb(slide, l_text, y, w_text, desc_h, desc, 9, color=text_color)

    # Link
    if url:
        _link_tb(slide, l_text, link_y, w_text, Inches(0.26),
                 "→ Mehr Informationen erhalten Sie hier", url)


# ══════════════════════════════════════════════════════════════════════════════
# Folienverwaltung
# ══════════════════════════════════════════════════════════════════════════════

def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in list(slide.placeholders):
        ph.element.getparent().remove(ph.element)
    return slide


def _copy_element_remapping_rels(shape_el, src_part, dst_part):
    """
    Tiefe Kopie eines Shape-XML-Elements; alle r:embed / r:link / r:id
    Attribute werden in neue Beziehungen der Ziel-Folie umgemappt.
    Verhindert kaputte Referenzen und den PowerPoint-Reparatur-Dialog.
    """
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    el = copy.deepcopy(shape_el)
    for attr in (f"{{{R}}}embed", f"{{{R}}}link", f"{{{R}}}id"):
        for node in el.iter():
            old_rid = node.get(attr)
            if not old_rid:
                continue
            if old_rid not in src_part.rels:
                node.attrib.pop(attr, None)
                continue
            rel = src_part.rels[old_rid]
            try:
                if rel.is_external:
                    new_rid = dst_part.relate_to(
                        rel.target_ref, rel.reltype, is_external=True)
                else:
                    new_rid = dst_part.relate_to(rel.target_part, rel.reltype)
                node.set(attr, new_rid)
            except Exception:
                node.attrib.pop(attr, None)
    return el


# ══════════════════════════════════════════════════════════════════════════════
# Inhaltsfolien (Artikel)
# ══════════════════════════════════════════════════════════════════════════════

def _content_slide(prs, articles, page_num, first_num, logo, year):
    if not articles:
        return
    slide = _new_slide(prs)
    _header(slide, logo)
    _footer(slide, logo, page_num, year)

    n = len(articles)
    total_h = CONTENT_BOT - CONTENT_TOP
    block_h = total_h / n

    for i, article in enumerate(articles):
        top = CONTENT_TOP + i * block_h
        _article_block(slide, article, first_num + i,
                       top, block_h,
                       MARGIN, TEXT_L_FULL, TEXT_W_FULL)
        if i < n - 1:
            # Trennlinie unterhalb des Autortexts (nicht überlappend)
            _divider(slide, top + block_h - Inches(0.02))


# ══════════════════════════════════════════════════════════════════════════════
# Event-Inhaltsfolien
# ══════════════════════════════════════════════════════════════════════════════

def _event_slide(prs, events, page_num, logo, year, show_more=False):
    """Erstellt eine Event-Folie mit bis zu 4 Events."""
    if not events:
        return
    slide = _new_slide(prs)
    _header(slide, logo)
    _footer(slide, logo, page_num, year)

    n       = len(events)
    total_h = CONTENT_BOT - CONTENT_TOP
    # "Mehr"-Block zählt als vollwertige Section – gleiche Höhe wie ein Event
    n_blocks = n + 1 if show_more else n
    block_h  = total_h / n_blocks

    for i, event in enumerate(events):
        top = CONTENT_TOP + i * block_h
        _event_block(slide, event, i + 1, top, block_h,
                     MARGIN, TEXT_L_FULL, TEXT_W_FULL)
        # Trennstrich nur wenn danach noch etwas folgt
        if i < n - 1 or show_more:
            _divider(slide, top + block_h - Inches(0.02))

    if show_more:
        top = CONTENT_TOP + n * block_h
        pad = Inches(0.12)
        y   = top + pad

        # Emoji statt Datum/Nummer
        _tb(slide, MARGIN, y, NUM_W, Inches(0.55),
            "📅", 20, color=IBM_BLUE, align=PP_ALIGN.CENTER)

        # Titel
        _tb(slide, TEXT_L_FULL, y, TEXT_W_FULL, Inches(0.50),
            "Weitere Events für diesen Zeitraum", 11, bold=True, color=NEAR_BLACK)
        y += Inches(0.56)

        # Beschreibung
        _tb(slide, TEXT_L_FULL, y, TEXT_W_FULL, Inches(0.26),
            "Für diesen Zeitraum sind noch weitere Events verfügbar.",
            9, color=MID_GRAY, italic=True)

        # Link
        link_y = top + block_h - Inches(0.28)
        _link_tb(slide, TEXT_L_FULL, link_y, TEXT_W_FULL, Inches(0.26),
                 "→ Alle Events auf der IBM Z Community Seite ansehen", EVENTS_URL)


# ══════════════════════════════════════════════════════════════════════════════
# Titelfolie (Cover)
# ══════════════════════════════════════════════════════════════════════════════

def _cover_slide(prs_slide, banner, logo,
                 month_name, year, issue_str, articles, events=None,
                 events_truncated=False):
    """
    Baut die Titelfolie komplett neu auf.
    events=None  → kein Events-Abschnitt in der Seitenleiste
    events=[]    → Events-Abschnitt mit "Keine Events" Hinweis
    events=[...] → Events-Abschnitt mit Event-Liste
    """
    slide = prs_slide

    KEEP = {"Grafik 4", "Grafik 82", "Picture 14", "Picture 46",
            "Rechteck 53", "Rechteck 80", "Textfeld 5"}

    to_del = []
    for shape in slide.shapes:
        if shape.name in KEEP:
            continue
        if shape.shape_type == 13:
            continue
        to_del.append(shape._element)
    for el in to_del:
        try:
            slide.shapes._spTree.remove(el)
        except Exception:
            pass

    # ── Seitenleiste ──────────────────────────────────────────────────────────
    sidebar_top = Inches(2.00)
    sidebar_h   = H - sidebar_top
    _rect(slide, Inches(0), sidebar_top, SIDE_W, sidebar_h, BAR_GRAY)

    # ── Issue-Nummer ──────────────────────────────────────────────────────────
    _tb(slide, Inches(0.20), Inches(1.55), Inches(2.0), Inches(0.34),
        f"Issue No. {issue_str}", 9, color=NEAR_BLACK)

    # ── Monat/Jahr-Box ────────────────────────────────────────────────────────
    _rect(slide, Inches(0.07), Inches(0.0), Inches(1.22), Inches(1.40),
          RGBColor(0x00, 0x2D, 0x9C))
    _tb(slide, Inches(0.10), Inches(0.12), Inches(1.10), Inches(0.60),
        month_name, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(0.10), Inches(0.68), Inches(1.10), Inches(0.46),
        str(year), 10, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Seitenleiste: Inhalt ──────────────────────────────────────────────────
    # Verfügbare Höhe: von 2.60" bis 10.80"
    ITEMS_TOP = Inches(2.60)
    ITEMS_BOT = Inches(10.80)
    available_h = ITEMS_BOT - ITEMS_TOP

    show_events = events is not None
    n_articles  = len(articles)
    n_events    = len(events) if events else 0

    # Platzvergabe
    if show_events:
        # Events-Abschnitt: Header + Einträge (oder Hinweis)
        evt_header_h = Inches(0.50)
        evt_empty_h  = Inches(0.36)
        evt_row_h    = Inches(0.46)
        evt_need_h   = evt_header_h + (n_events * evt_row_h if n_events else evt_empty_h)
        events_h     = min(evt_need_h, available_h * 0.38)
        gap_h        = Inches(0.20)
        articles_h   = available_h - events_h - gap_h
    else:
        articles_h   = available_h
        events_h     = 0
        gap_h        = 0

    article_row_h = max(
        Inches(0.30),
        min(Inches(0.65), articles_h / n_articles if n_articles else Inches(0.65)),
    )

    # ── THEMEN-Überschrift ────────────────────────────────────────────────────
    _tb(slide, SIDE_COL_L, Inches(2.10), SIDE_COL_W, Inches(0.44),
        "THEMEN", 12, bold=True, color=NEAR_BLACK)

    # ── Artikel-Liste ─────────────────────────────────────────────────────────
    # Zahl-Spalte: wrap=False verhindert Umbruch bei zweistelligen Zahlen
    NUM_COL_W  = Inches(0.46)
    NUM_COL_OFF = Inches(0.50)
    y = ITEMS_TOP
    for i, art in enumerate(articles):
        _tb(slide, SIDE_COL_L, y, NUM_COL_W, article_row_h,
            str(i + 1), 9, bold=True, color=IBM_BLUE, align=PP_ALIGN.CENTER,
            wrap=False)
        _tb(slide, SIDE_COL_L + NUM_COL_OFF, y,
            SIDE_COL_W - NUM_COL_OFF, article_row_h,
            art["title"], 7, color=NEAR_BLACK, wrap=True)
        y += article_row_h

    # ── EVENTS-Abschnitt ──────────────────────────────────────────────────────
    if show_events:
        y += gap_h
        # Trennlinie
        _rect(slide, SIDE_COL_L, y - Inches(0.08),
              SIDE_COL_W, Inches(0.012), DIV_GRAY)

        _tb(slide, SIDE_COL_L, y, SIDE_COL_W, Inches(0.40),
            "EVENTS", 12, bold=True, color=NEAR_BLACK)
        y += Inches(0.44)

        if n_events == 0:
            _tb(slide, SIDE_COL_L, y, SIDE_COL_W, Inches(0.36),
                "Keine Events für diesen Zeitraum",
                7, color=MID_GRAY, italic=True)
        else:
            remaining_h = ITEMS_BOT - y
            actual_row_h = max(
                Inches(0.32),
                min(Inches(0.50), remaining_h / n_events),
            )
            for evt in events:
                pub = evt.get("event_date") if isinstance(evt, dict) else getattr(evt, "event_date", None)
                title = evt.get("title") if isinstance(evt, dict) else evt.title
                date_str = (f"{pub.day}.{pub.month:02d}.  " if pub else "")
                _tb(slide, SIDE_COL_L, y, SIDE_COL_W, actual_row_h,
                    f"{date_str}{title}", 7, color=NEAR_BLACK, wrap=True)
                y += actual_row_h
            if events_truncated:
                _tb(slide, SIDE_COL_L, y, SIDE_COL_W, Inches(0.30),
                    "und vieles mehr ...", 7, color=MID_GRAY, italic=True)

    # ── Seitenzahl ────────────────────────────────────────────────────────────
    _tb(slide, Inches(3.90), Inches(11.36), Inches(0.50), Inches(0.28),
        "1", 12, color=NEAR_BLACK, align=PP_ALIGN.CENTER)

    # ── Artikel 1 auf der rechten Seite ───────────────────────────────────────
    if articles:
        _article_block(slide, articles[0], 1,
                       Inches(2.05), Inches(8.70),
                       COVER_R_L, COVER_R_L + NUM_W + Inches(0.06),
                       COVER_R_W - NUM_W - Inches(0.06))


# ══════════════════════════════════════════════════════════════════════════════
# Folie löschen
# ══════════════════════════════════════════════════════════════════════════════

def _delete_slide(prs, index):
    slides = prs.slides
    slide  = list(slides)[index]
    lst    = slides._sldIdLst
    sldId_elem = list(lst)[index]
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    rid  = sldId_elem.get(f"{{{r_ns}}}id")
    prs.part.drop_rel(rid)
    lst.remove(sldId_elem)


# ══════════════════════════════════════════════════════════════════════════════
# Hauptfunktion
# ══════════════════════════════════════════════════════════════════════════════

def build_newsletter(articles, month, year, issue_number,
                     events=None, output_filename=None, events_truncated=False):
    """
    Erstellt den Newsletter als PPTX.

    articles      – Liste von Artikel-Dicts (title, author, url, published,
                    summary, image_url, image_bytes)
    month / year  – Monat und Jahr der Ausgabe
    issue_number  – Ausgabe-Nummer als String
    events        – None (kein Events-Abschnitt) oder Liste von Event-Dicts
                    (title, event_date, time_str, location, description, url)
    """
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    month_name = GERMAN_MONTHS.get(month, str(month))
    fname = output_filename or f"IBM_Z_Newsletter_{month_name}_{year}.pptx"
    out   = os.path.join(OUTPUT_DIR, fname)
    shutil.copy2(TEMPLATE_FILE, out)

    prs    = Presentation(out)
    banner, logo = _extract_images(prs)

    # ── Inhalt der Abschlussfolie aus Template sichern ────────────────────────
    # Bilder werden als Bytes extrahiert (nicht als XML-Referenz), damit die
    # rId-Beziehungen beim Neuaufbau der Folie korrekt gesetzt werden.
    closing_items = []   # list of dicts: {type, ...}
    tpl_closing = list(prs.slides)[-1]

    # Fallback: Logo aus der Closing-Slide extrahieren falls noch nicht gefunden
    if logo is None:
        for shape in tpl_closing.shapes:
            if (shape.shape_type == 13
                    and shape.top < Inches(1.0)
                    and shape.left > Inches(6.0)):
                try:
                    logo = shape.image.blob
                except Exception:
                    pass
                break

    for shape in tpl_closing.shapes:
        # Haupttitel überspringen – wird neu gerendert
        if shape.has_text_frame and shape.text_frame.text.strip() == "ZUSÄTZLICHE INFORMATIONEN":
            continue
        # Header-Bereich überspringen – wird per _header() neu gesetzt
        if shape.top < Inches(1.0):
            continue
        if shape.shape_type == 13:
            # Bild: Bytes + Geometrie sichern, NICHT das XML (würde rId verlieren)
            try:
                closing_items.append({
                    "type": "image",
                    "bytes": shape.image.blob,
                    "left": shape.left, "top": shape.top,
                    "width": shape.width, "height": shape.height,
                })
            except Exception:
                pass
        else:
            # Text/Vektorelement: rIds korrekt remappen
            closing_items.append({
                "type": "shape",
                "element": shape._element,   # Original; wird beim Rebuild neu gemappt
                "src_part": tpl_closing.part,
            })

    # ── Titelfolie aufbauen ───────────────────────────────────────────────────
    _cover_slide(prs.slides[0], banner, logo,
                 month_name, year, issue_number or "?",
                 articles, events, events_truncated=events_truncated)

    # ── ALLE Template-Folien außer Cover löschen ──────────────────────────────
    # Verhindert Namenskonflikte (z.B. slide7.xml) beim Hinzufügen neuer Folien
    while len(prs.slides) > 1:
        _delete_slide(prs, 1)

    # ── Artikel-Inhaltsfolien (Artikel 2+ in 2er-Gruppen) ────────────────────
    remaining = articles[1:]
    groups    = [remaining[i:i + 2] for i in range(0, len(remaining), 2)] if remaining else []

    for page_idx, group in enumerate(groups):
        first_num = 2 + page_idx * 2
        _content_slide(prs, group, page_idx + 2, first_num, logo, year)

    # ── Event-Folien (4 Events pro Folie) ────────────────────────────────────
    n_article_slides = len(groups)
    if events:
        events_per_slide = 4
        event_groups = [events[i:i + events_per_slide]
                        for i in range(0, len(events), events_per_slide)]
        for eg_idx, eg in enumerate(event_groups):
            page_num  = 2 + n_article_slides + eg_idx
            is_last   = eg_idx == len(event_groups) - 1
            _event_slide(prs, eg, page_num, logo, year,
                         show_more=events_truncated and is_last)

    # ── Abschlussfolie neu aufbauen (immer letzte Folie) ─────────────────────
    closing = _new_slide(prs)
    _header(closing, logo)
    # Template-Inhalte einfügen – Bilder per add_picture (korrekte rIds)
    for item in closing_items:
        try:
            if item["type"] == "image":
                closing.shapes.add_picture(
                    io.BytesIO(item["bytes"]),
                    item["left"], item["top"], item["width"], item["height"],
                )
            else:
                el = _copy_element_remapping_rels(
                    item["element"], item["src_part"], closing.part)
                closing.shapes._spTree.append(el)
        except Exception:
            pass
    # Titel
    _tb(closing, TEXT_L_FULL, CONTENT_TOP, TEXT_W_FULL, Inches(0.45),
        "Zusätzliche Informationen", 16, bold=True, color=NEAR_BLACK)
    # "Besuchen Sie uns"-Block korrekt positionieren
    for shape in closing.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip().startswith("Besuchen Sie uns"):
            shape.left  = TEXT_L_FULL
            shape.width = TEXT_W_FULL
            shape.top   = CONTENT_TOP + Inches(0.55)
    _footer(closing, logo, len(prs.slides), year)

    prs.save(out)
    return out
